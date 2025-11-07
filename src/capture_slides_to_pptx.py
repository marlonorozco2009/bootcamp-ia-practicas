# capture_slides_to_pptx.py
# Agente: captura cada slide de una presentación pública de Google Slides (pubembed)
# y arma un PPTX local llamado "presentacion_capturada.pptx".
#
# Nota: si la presentación necesita login no funcionará. Ajusta MAX_SLIDES si hace falta.

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import hashlib, os, time, io

# URL pública (la que nos diste). Hecho un pequeño "clean" al parámetro amp.
URL = "https://docs.google.com/presentation/d/e/2PACX-1vTvp5dUYrT9ZBcUpcclAVsL6tUMCXpGgjOpMbIosLyH7etVbEMe5dBjwkKyF_Hq_pi0QX8oVWhRGVfJ/pubembed?amp;usp=embed_facebook#slide=id.g38e3d3008be_4_166"

OUT_DIR = "slides_images"
OUT_PPTX = "presentacion_capturada.pptx"
MAX_SLIDES = 300            # corte por seguridad
WAIT_AFTER_NAV = 0.8        # segundos (ajusta si tu conexión es lenta)

def md5_bytes(b):
    return hashlib.md5(b).hexdigest()

def ensure_dir(d):
    if not os.path.exists(d):
        os.makedirs(d)

def image_bytes_of_element(page, selector=None):
    # si se pasa selector, intenta screenshot del elemento; si no, screenshot de viewport
    if selector:
        el = page.query_selector(selector)
        if el:
            return el.screenshot(type="png")
    return page.screenshot(type="png", full_page=False)

def main():
    ensure_dir(OUT_DIR)
    images = []
    prev_hash = None

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True, args=["--window-size=1920,1080"])
        context = browser.new_context(viewport={"width":1920, "height":1080})
        page = context.new_page()
        print("[*] Abriendo URL...")
        page.goto(URL, wait_until="networkidle", timeout=60000)
        time.sleep(1.0)  # deja que todo renderice

        # Intenta identificar un contenedor de slide típico para recortar mejor
        # Si no existe, hará screenshot del viewport.
        possible_selectors = [
            'iframe',                         # pubembed a menudo tiene iframe
            '.punch-viewer',                  # componente interno (si aparece)
            'div[role="presentation"]',       # heurística
            '#canvas',                        # heurística
        ]
        chosen_selector = None
        for sel in possible_selectors:
            try:
                if page.query_selector(sel):
                    chosen_selector = sel
                    break
            except Exception:
                pass

        if chosen_selector:
            print(f"[*] Usando selector para captura: {chosen_selector}")
        else:
            print("[*] No se encontró selector específico. Se usará viewport.")

        # Ir a la primera diapositiva (si el fragment lo pone, ok). Luego iterar con ArrowRight.
        # Guardar cada captura; si detecta repetición (hash idéntico) dos veces seguidas, se corta.
        for i in range(MAX_SLIDES):
            # Espera pequeño e intenta screenshot del elemento
            time.sleep(0.5)
            try:
                img_bytes = image_bytes_of_element(page, chosen_selector)
            except Exception as e:
                print("[!] Error haciendo screenshot:", e)
                img_bytes = page.screenshot(type="png")

            h = md5_bytes(img_bytes)
            if h == prev_hash:
                # Si la imagen no cambió al avanzar, probamos una vez más y si repite, cortamos.
                # (evita cortar inmediatamente si la primera captura se repite por render)
                print(f"[*] Imagen igual a la anterior (posible fin). Iter {i}. Cortando.")
                break

            filename = os.path.join(OUT_DIR, f"slide_{i+1:03d}.png")
            with open(filename, "wb") as f:
                f.write(img_bytes)
            images.append(filename)
            print(f"[+] Capturada slide {i+1} -> {filename}")

            prev_hash = h

            # avanzar a la siguiente slide
            try:
                page.keyboard.press("ArrowRight")
            except Exception:
                # fallback: click en zona central para forzar avance
                try:
                    page.mouse.click(960, 540)
                except Exception:
                    pass
            time.sleep(WAIT_AFTER_NAV)

        print(f"[*] Capturadas {len(images)} imágenes. Construyendo PPTX...")

        # Construir PPTX
        prs = Presentation()
        # usamos tamaño por defecto (10 x 7.5 inches). Vamos a ajustar la imagen para ocupar toda la slide.
        slide_w = prs.slide_width
        slide_h = prs.slide_height

        for img_path in images:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # layout en blanco
            # abrir imagen con Pillow para conocer tamaño y mantener aspecto
            with Image.open(img_path) as im:
                img_w_px, img_h_px = im.size
            # Ajuste simple: intentar llenar toda la slide preservando aspecto
            # convertimos px a inches aproximando 96 DPI (esto es aproximado). Mejor: usar dimensiones EMU directas.
            # Pero python-pptx acepta width/height en EMU o usando Inches. Vamos a calcular una escala en base al ratio.
            ratio_slide = slide_w / slide_h
            ratio_img = img_w_px / img_h_px

            if ratio_img > ratio_slide:
                # imagen más ancha -> ajustar ancho a slide_w
                pic = slide.shapes.add_picture(img_path, 0, 0, width=slide_w)
                # centrar verticalmente
                top = int((slide_h - pic.height) // 2)
                pic.top = top
            else:
                # imagen más alta -> ajustar altura a slide_h
                pic = slide.shapes.add_picture(img_path, 0, 0, height=slide_h)
                left = int((slide_w - pic.width) // 2)
                pic.left = left

        prs.save(OUT_PPTX)
        print(f"[✓] PPTX guardado como: {OUT_PPTX}")
        browser.close()

if __name__ == "__main__":
    main()
